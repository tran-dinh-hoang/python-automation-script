import os
from pptx import Presentation
import openpyxl
from openpyxl import Workbook

# Function to extract hyperlinks from a PowerPoint file
def extract_hyperlinks_from_pptx(pptx_file):
    presentation = Presentation(pptx_file)
    hyperlinks = []

    # Iterate through all slides and shapes to find hyperlinks
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:  # Check if the shape has text
                if shape.hyperlink and shape.hyperlink.address:  # Check if the shape has a hyperlink
                    title = shape.text
                    hyperlink = shape.hyperlink.address
                    hyperlinks.append((title.strip(), hyperlink))
    
    return hyperlinks

# Function to write hyperlinks to an Excel file in the output folder
def write_to_excel(file_name, data, output_folder):
    workbook = Workbook()
    sheet = workbook.active

    # First row: file name
    sheet["A1"] = "File Name"
    sheet["B1"] = file_name

    # Second row: headers
    sheet["A2"] = "Title"
    sheet["B2"] = "Hyperlink"

    # From the third row: data (title and hyperlink)
    for idx, (title, hyperlink) in enumerate(data, start=3):
        sheet[f"A{idx}"] = title
        sheet[f"B{idx}"] = hyperlink

    # Ensure output folder exists
    os.makedirs(output_folder, exist_ok=True)

    # Save the Excel file in the output folder
    output_file_name = os.path.basename(file_name).replace(".pptx", ".xlsx")
    output_file_path = os.path.join(output_folder, output_file_name)
    workbook.save(output_file_path)

    print(f"Data written to {output_file_path}")

# Main function
def main():
    # Folder where the PowerPoint file is located
    input_folder = "input"
    
    # Folder where the Excel file will be saved
    output_folder = "output"

    # List all files in the folder and select the first PowerPoint file
    pptx_files = [f for f in os.listdir(input_folder) if f.endswith(".pptx")]

    if pptx_files:
        pptx_file_path = os.path.join(input_folder, pptx_files[0])  # Get the first PowerPoint file in the folder
        
        # Extract hyperlinks
        hyperlinks = extract_hyperlinks_from_pptx(pptx_file_path)

        # Write data to Excel
        if hyperlinks:
            write_to_excel(pptx_file_path, hyperlinks, output_folder)
        else:
            print("No hyperlinks found in the PowerPoint.")
    else:
        print("No PowerPoint files found in the input folder.")

# Run the script
if __name__ == "__main__":
    main()
